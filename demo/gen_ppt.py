"""
Generate KavachVoice SIH 2026 presentation (5 slides).
Usage: python demo/gen_ppt.py
Output: demo/KavachVoice_SIH2026.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "KavachVoice_SIH2026.pptx"

# Brand colors
BLUE_DARK  = RGBColor(0x15, 0x65, 0xC0)   # #1565C0
BLUE_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)   # #E3F2FD
RED        = RGBColor(0xD5, 0x00, 0x00)   # #D50000
GREEN      = RGBColor(0x1B, 0x5E, 0x20)   # #1B5E20
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x21, 0x21, 0x21)
GREY       = RGBColor(0x75, 0x75, 0x75)

W = Inches(13.333)   # widescreen 16:9
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor = None, line: RGBColor = None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=20, color=DARK_TEXT, bullet="▸ "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def header_bar(slide, title, subtitle=None):
    """Blue top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.25), fill=BLUE_DARK)
    add_text(slide, title,
             Inches(0.4), Inches(0.1), Inches(10), Inches(0.75),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.82), Inches(10), Inches(0.35),
                 font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)


def footer(slide):
    add_text(slide, "KavachVoice  |  SIH 2026  |  Problem SIH26104  |  Theme: Blockchain & Cybersecurity  |  Sponsor: AICTE",
             Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.35),
             font_size=10, color=GREY, align=PP_ALIGN.CENTER)


# ─── Slide 1: Title ──────────────────────────────────────────────────────────
def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BLUE_DARK)

    # White card
    add_rect(sl, Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.2), fill=WHITE)

    add_text(sl, "KavachVoice",
             Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
             font_size=60, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    add_text(sl, "Three-Layer Voice Fraud Shield for India",
             Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
             font_size=26, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text(sl, "Proactive Perturbation  ·  Real-Time Call Guard  ·  AI Forensics",
             Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.5),
             font_size=18, color=GREY, align=PP_ALIGN.CENTER)

    add_text(sl, "SIH 2026  ·  Problem SIH26104  ·  Theme: Blockchain & Cybersecurity  ·  Sponsor: AICTE",
             Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.4),
             font_size=14, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)


# ─── Slide 2: Problem ────────────────────────────────────────────────────────
def slide_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    header_bar(sl, "The Problem", subtitle="Voice fraud is India's fastest-growing cybercrime vector")

    # Big stat
    add_rect(sl, Inches(0.4), Inches(1.4), Inches(3.8), Inches(2.0), fill=RGBColor(0xFF, 0xEB, 0xEE))
    add_text(sl, "₹11,333 Cr",
             Inches(0.45), Inches(1.5), Inches(3.7), Inches(0.9),
             font_size=38, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(sl, "lost to cyber fraud in India (2023–24)\nsource: MHA Annual Report",
             Inches(0.45), Inches(2.35), Inches(3.7), Inches(0.9),
             font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(4.5), Inches(1.4), Inches(3.8), Inches(2.0), fill=RGBColor(0xFF, 0xEB, 0xEE))
    add_text(sl, "77%",
             Inches(4.55), Inches(1.5), Inches(3.7), Inches(0.9),
             font_size=38, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(sl, "of vishing victims comply\nbecause the voice sounds real",
             Inches(4.55), Inches(2.35), Inches(3.7), Inches(0.9),
             font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(8.6), Inches(1.4), Inches(4.3), Inches(2.0), fill=RGBColor(0xFF, 0xEB, 0xEE))
    add_text(sl, "AWS & Azure",
             Inches(8.65), Inches(1.5), Inches(4.2), Inches(0.7),
             font_size=28, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(sl, "both retired their voice fraud\ndetection APIs in 2024 —\nno commercial solution exists",
             Inches(8.65), Inches(2.15), Inches(4.2), Inches(1.1),
             font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text(sl, "Attack chain: Attacker clones victim's family member voice with ElevenLabs/Coqui (free, 30-second sample) → calls during UPI transaction → victim transfers money before verifying.",
             Inches(0.4), Inches(3.6), Inches(12.5), Inches(0.8),
             font_size=16, color=DARK_TEXT)

    add_bullet_box(sl,
        ["No existing app defends against voice cloning proactively",
         "Detection-only tools can't stop a transfer already in progress",
         "No forensic evidence trail for I4C / CERT-In reporting"],
        Inches(0.4), Inches(4.5), Inches(12.5), Inches(1.8),
        font_size=17, color=DARK_TEXT)

    footer(sl)


# ─── Slide 3: Solution ───────────────────────────────────────────────────────
def slide_solution(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    header_bar(sl, "KavachVoice — Three-Layer Defense", subtitle="The only solution that is proactive, real-time, AND evidence-grade")

    # Layer boxes
    layer_tops = [Inches(1.4), Inches(3.05), Inches(4.7)]
    layer_colors = [RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0xFF, 0xF3, 0xE0)]
    layer_accent = [GREEN, BLUE_DARK, RGBColor(0xE6, 0x51, 0x00)]
    layer_titles = ["Layer 1 — VoiceArmor  (Proactive)", "Layer 2 — CallGuard  (Real-Time)", "Layer 3 — VoiceID SDK  (Forensic)"]
    layer_bodies = [
        "Universal Adversarial Perturbation applied to PCM audio in <50 ms. Disrupts harmonic structure of voice data captured by cloning tools — poisoned audio produces degraded synthetic output. No user action needed.",
        "Accessibility Service detects UPI app foreground during active phone call. Keyword spotting (OTP / जल्दी / बताइए / urgency terms) triggers yellow → orange → red overlay. TYPE_ACCESSIBILITY_OVERLAY — no root required.",
        "RawNet2 + ECAPA-TDNN run in parallel via asyncio (FastAPI). Sub-800 ms inference. Returns SYNTHETIC / GENUINE / UNCERTAIN verdict. SHA-256 audio hash + confidence scores → forensic PDF → I4C submission.",
    ]

    for i in range(3):
        add_rect(sl, Inches(0.3), layer_tops[i], Inches(12.7), Inches(1.5), fill=layer_colors[i])
        add_rect(sl, Inches(0.3), layer_tops[i], Inches(0.15), Inches(1.5), fill=layer_accent[i])
        add_text(sl, layer_titles[i],
                 Inches(0.6), layer_tops[i] + Pt(8), Inches(12.1), Inches(0.45),
                 font_size=18, bold=True, color=layer_accent[i])
        add_text(sl, layer_bodies[i],
                 Inches(0.6), layer_tops[i] + Inches(0.5), Inches(12.1), Inches(0.9),
                 font_size=15, color=DARK_TEXT)

    footer(sl)


# ─── Slide 4: Demo Flow ──────────────────────────────────────────────────────
def slide_demo(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    header_bar(sl, "Live Demo — 4 Moments  (<5 min)", subtitle="Every moment is independently executable with a fallback")

    moments = [
        ("1", "Clone Attack Detected",     "POST clone_raw.wav (ElevenLabs) → SYNTHETIC verdict · confidence 100% · latency <10 ms",  BLUE_DARK),
        ("2", "VoiceArmor Spectrograms",   "Before: smooth harmonics. After UAP: disrupted structure. Cloning algorithm output degrades to noise.", GREEN),
        ("3", "UPI Block — Red Overlay",   "Call active + PhonePe/GPay open + 'OTP bataiye' → red TYPE_ACCESSIBILITY_OVERLAY in <2 s", RED),
        ("4", "Forensic PDF + I4C Button", "SYNTHETIC verdict → auto PDF (SHA-256, scores, session ID, timestamp) → 'Submit to I4C [DEMO MODE]'", RGBColor(0xE6, 0x51, 0x00)),
    ]

    for i, (num, title, body, color) in enumerate(moments):
        top = Inches(1.4) + i * Inches(1.35)
        add_rect(sl, Inches(0.3), top, Inches(0.7), Inches(1.1), fill=color)
        add_text(sl, num, Inches(0.3), top + Pt(10), Inches(0.7), Inches(0.7),
                 font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(1.1), top, Inches(11.9), Inches(1.1), fill=BLUE_LIGHT if i % 2 == 0 else RGBColor(0xFA, 0xFA, 0xFA))
        add_text(sl, title, Inches(1.2), top + Pt(6), Inches(11.7), Inches(0.45),
                 font_size=18, bold=True, color=color)
        add_text(sl, body, Inches(1.2), top + Inches(0.5), Inches(11.7), Inches(0.55),
                 font_size=14, color=DARK_TEXT)

    footer(sl)


# ─── Slide 5: Why KavachVoice Wins ──────────────────────────────────────────
def slide_why(prs):
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    header_bar(sl, "Why KavachVoice Wins", subtitle="Market gap + unique 3-layer architecture no competitor has")

    # Comparison table header
    cols = [Inches(0.3), Inches(4.5), Inches(7.0), Inches(9.5), Inches(11.5)]
    col_w = [Inches(4.1), Inches(2.4), Inches(2.4), Inches(1.9), Inches(1.7)]
    row_h = Inches(0.48)
    rows_top = Inches(1.45)

    headers = ["Capability", "KavachVoice", "AWS (retired 2024)", "Azure (retired 2024)", "Truecaller"]
    row_data = [
        ["Proactive voice poisoning",    "YES",  "NO",  "NO",  "NO"],
        ["Real-time call intervention",  "YES",  "NO",  "NO",  "Partial"],
        ["No root / system perms needed","YES",  "N/A", "N/A", "YES"],
        ["Forensic PDF + evidence chain","YES",  "NO",  "NO",  "NO"],
        ["I4C submission ready",         "YES",  "NO",  "NO",  "NO"],
        ["Works offline (Android)",      "YES",  "NO",  "NO",  "Partial"],
    ]

    def cell_color(val):
        if val == "YES": return RGBColor(0xC8, 0xE6, 0xC9)
        if val == "NO":  return RGBColor(0xFF, 0xCD, 0xD2)
        return RGBColor(0xFF, 0xF9, 0xC4)

    # Header row
    for j, h in enumerate(headers):
        bg = BLUE_DARK if j == 0 else (BLUE_DARK if j == 1 else RGBColor(0x61, 0x61, 0x61))
        add_rect(sl, cols[j], rows_top, col_w[j], row_h, fill=bg)
        add_text(sl, h, cols[j] + Pt(4), rows_top + Pt(4), col_w[j] - Pt(8), row_h,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for r, row in enumerate(row_data):
        top = rows_top + (r + 1) * row_h
        for j, val in enumerate(row):
            bg = RGBColor(0xF5, 0xF5, 0xF5) if j == 0 else cell_color(val)
            add_rect(sl, cols[j], top, col_w[j], row_h, fill=bg)
            fc = DARK_TEXT if j == 0 else (GREEN if val == "YES" else (RED if val == "NO" else DARK_TEXT))
            add_text(sl, val, cols[j] + Pt(4), top + Pt(4), col_w[j] - Pt(8), row_h,
                     font_size=13, bold=(j == 0), color=fc,
                     align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    add_text(sl, '"AWS Rekognition Voice Liveness and Azure Speaker Recognition were both discontinued in 2024.\nKavachVoice fills the gap with a fully on-device, open-source, evidence-grade solution."',
             Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.85),
             font_size=15, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    footer(sl)


def main():
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_demo(prs)
    slide_why(prs)
    prs.save(str(OUT))
    print(f"Saved: {OUT}  ({OUT.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
